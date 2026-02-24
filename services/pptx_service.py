from pptx import Presentation
from datetime import datetime
import os

class PPTXService:
    @staticmethod
    def extract_slides(pptx_path):
        prs = Presentation(pptx_path)
        slides_content = []
        
        for idx, slide in enumerate(prs.slides):
            content = []
            
            # Extract title
            if slide.shapes.title:
                content.append(f"Title: {slide.shapes.title.text}")
            
            # Process all shapes
            for shape in slide.shapes:
                # Extract text
                if hasattr(shape, "text") and shape.text.strip():
                    content.append(shape.text)
                
                # Extract table data
                if shape.shape_type == 19:  # Table
                    table_data = PPTXService._extract_table(shape)
                    if table_data:
                        content.append(table_data)
                
                # Describe charts
                if shape.shape_type == 3:  # Chart
                    if hasattr(shape, "chart"):
                        chart_info = PPTXService._describe_chart(shape)
                        content.append(chart_info)
                    else:
                        content.append("[Chart/Graph present]")
                
                # Note other visual elements
                if shape.shape_type == 13:  # Picture
                    content.append("[Image/Picture present]")
            
            # Ensure every slide has content
            if not content:
                content.append(f"Slide {idx + 1} - Visual content")
            
            slides_content.append({
                'index': idx,
                'content': '\n'.join(content)
            })
        
        return slides_content, prs
    
    @staticmethod
    def _extract_table(shape):
        try:
            table = shape.table
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append(" | ".join(cells))
            return "TABLE:\n" + "\n".join(rows)
        except:
            return "[Table present]"
    
    @staticmethod
    def _describe_chart(shape):
        try:
            chart = shape.chart
            chart_type = str(chart.chart_type)
            return f"[Chart: {chart_type}]"
        except:
            return "[Chart present]"
    
    @staticmethod
    def add_talking_points(prs, slide_analyses):
        for analysis in slide_analyses:
            slide_idx = analysis['slide_index']
            if slide_idx < len(prs.slides):
                slide = prs.slides[slide_idx]
                
                # Add talking points to notes
                if not slide.has_notes_slide:
                    slide.notes_slide
                
                notes_text = ""
                
                # Talking Points
                notes_text += "TALKING POINTS:\n"
                for i, point in enumerate(analysis['talking_points'], 1):
                    notes_text += f"{i}. {point}\n"
                
                # Action Items
                notes_text += "\nACTION ITEMS:\n"
                real_items = [item for item in analysis['action_items'] if 'none identified' not in item.lower()]
                if real_items:
                    for item in real_items:
                        notes_text += f"- {item}\n"
                else:
                    notes_text += "- None identified for this slide\n"
                
                # Q&A
                notes_text += "\nANTICIPATED QUESTIONS:\n"
                if analysis['questions']:
                    for question in analysis['questions']:
                        notes_text += f"{question}\n\n"
                else:
                    notes_text += "- None identified for this slide\n"
                
                slide.notes_slide.notes_text_frame.text = notes_text
        
        return prs
    
    @staticmethod
    def save_presentation(prs, customer_name, output_folder):
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{customer_name}_MBR_{timestamp}.pptx"
        filepath = os.path.join(output_folder, filename)
        prs.save(filepath)
        return filename
    
    @staticmethod
    def save_action_items(slide_analyses, customer_name, output_folder):
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{customer_name}_ActionItems_{timestamp}.md"
        filepath = os.path.join(output_folder, filename)
        
        with open(filepath, 'w') as f:
            f.write(f"# Action Items - {customer_name} MBR\n\n")
            
            has_actions = False
            for analysis in slide_analyses:
                items = analysis['action_items']
                # Filter out "None identified" messages
                real_items = [item for item in items if 'none identified' not in item.lower()]
                
                if real_items:
                    has_actions = True
                    f.write(f"## Slide {analysis['slide_index'] + 1}\n")
                    for item in real_items:
                        f.write(f"- {item}\n")
                    f.write("\n")
            
            if not has_actions:
                f.write("No action items identified across all slides.\n")
        
        return filename
    
    @staticmethod
    def save_qa(slide_analyses, customer_name, output_folder):
        timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        filename = f"{customer_name}_QA_{timestamp}.md"
        filepath = os.path.join(output_folder, filename)
        
        with open(filepath, 'w') as f:
            f.write(f"# Q&A - {customer_name} MBR\n\n")
            for analysis in slide_analyses:
                if analysis['questions']:
                    f.write(f"## Slide {analysis['slide_index'] + 1}\n")
                    for question in analysis['questions']:
                        f.write(f"{question}\n\n")
        
        return filename
